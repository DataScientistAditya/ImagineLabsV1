from pptx import Presentation
from pptx.util import Inches
from .Fetchdocs import FetchDocuments
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.util import Pt
import nltk
import os
import datetime
import string






# prs = Presentation()
# Layout = prs.slide_layouts[0] 
# first_slide = prs.slides.add_slide(Layout) # Adding first slide
# first_slide.shapes.title.text = title
# first_slide.placeholders[1].text = "Created by ImagineLabs"
# left = top = Inches(0)
# img_path = imagelinks[0]
# pic = first_slide.shapes.add_picture(img_path, left, top, width=prs.slide_width, height=prs.slide_height)

# # This moves it to the background
# first_slide.shapes._spTree.remove(pic._element)
# first_slide.shapes._spTree.insert(2, pic._element)
# Current_Time = datetime.datetime.now()
# date_time = Current_Time.strftime("%m/%d/%Y, %H:%M:%S")
# Splitdate = date_time.split(" ")
# Endswith = "".join(Splitdate[1].split(":"))
# FileName = "{0}{1}{2}".format(title,Endswith,".pptx")
# prs.save(FileName)
# os.startfile(FileName)



prs = Presentation()

class CreateDocument():
    
    def __init__(self,pptcontains=[],imagelist=[],title=None,Query=None):
        self.pptcontains = pptcontains
        self.imagelist =imagelist
        self.title = title
        Current_Time = datetime.datetime.now()
        date_time = Current_Time.strftime("%m/%d/%Y, %H:%M:%S")
        Splitdate = date_time.split(" ")
        Endswith = "".join(Splitdate[1].split(":"))
        FileName = "{0}{1}{2}".format(Query,Endswith,".pptx")
        self.File = FileName
        
    def CreateFirstSlide(self):
        Layout = prs.slide_layouts[0] 
        first_slide = prs.slides.add_slide(Layout)
        first_slide.shapes.title.text = self.title
        subtitle = first_slide.placeholders[1]
        subtitle.text = "Created by ImagineLabs"
    
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
        left = top = Inches(0)
        img_path = self.imagelist[0]
        pic = first_slide.shapes.add_picture(img_path, left, top, width=prs.slide_width, height=prs.slide_height)

        first_slide.shapes._spTree.remove(pic._element)
        first_slide.shapes._spTree.insert(2, pic._element)
        
    def remove_punctuation(self,text):
        for punctuation in string.punctuation:
            text = text.replace(punctuation, '')
        return text
    
    def slide_maker(self,head,d, slide,layout):
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = head 
        shape = slide.shapes
        body_shape = shape.placeholders[1]
        textframe = body_shape.text_frame
        textframe.clear()   # not necessary for newly-created shape
        textframe.vertical_anchor = MSO_ANCHOR.TOP
        textframe.word_wrap = True
        textframe.margin_top = 0
        textframe.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        for text in nltk.sent_tokenize(str(d)):
            body = self.remove_punctuation(str(text))
    #            print(body)
            para = str(body)
    #         print('para is: '+ para)
            paragraph = textframe.add_paragraph()

            run = paragraph.add_run()
            run.text = para
            run.alignment = PP_ALIGN.LEFT
            run.level = 1
            font = run.font
            font.name = 'Calibri'
            font.size = Pt(12)
            font.bold = False
            font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
            prs.save(self.File)
            
            
    def CreatePPT(self):
        for idx, des in enumerate(self.pptcontains):

            if len(des)<100:
                head = self.remove_punctuation(str(des))
                head_idx= idx

                for d in self.pptcontains[idx+1:idx+2]:
                        slide = 'slide_'+str(idx)
                        layout = 'layout_'+str(idx)
                        self.slide_maker(head,d,slide,layout)
                        break
                    
        return self.File
        # os.startfile(self.File)
        # X.save(self.File)
        # os.startfile(self.File)
        
        
        

    



